"""Script for converting paragraphs from .docx file to separate slides in .pptx file.

Keeps formating of text like: bold, italic, color, font, size and hightlight.
"""

from pathlib import Path
from tkinter import Tk, filedialog

from docx import Document
from docx.enum.text import WD_COLOR
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches


class Pptx:
    # Mapping based on https://docs.microsoft.com/en-us/office/vba/api/Word.WdColor
    WD_COLOR_TO_RGB_STR = {
        WD_COLOR.BLACK: '000000',
        WD_COLOR.BLUE: '0000FF',
        WD_COLOR.BRIGHT_GREEN: '00FF00',
        WD_COLOR.DARK_BLUE: '000080',
        WD_COLOR.DARK_RED: '800000',
        WD_COLOR.DARK_YELLOW: '808000',
        WD_COLOR.GRAY_25: 'C0C0C0',
        WD_COLOR.GRAY_50: '808080',
        WD_COLOR.GREEN: '008000',
        WD_COLOR.PINK: 'FF00FF',
        WD_COLOR.RED: 'FF0000',
        WD_COLOR.TEAL: '008080',
        WD_COLOR.TURQUOISE: '00FFFF',
        WD_COLOR.VIOLET: '800080',
        WD_COLOR.WHITE: 'FFFFFF',
        WD_COLOR.YELLOW: 'FFFF00',
    }

    def __init__(self) -> None:
        self.pptx = Presentation()  # I have exception thrown when inheriting from this class
        super().__init__()

    def new_blank_slide(self):
        pptx_slide = self.pptx.slides.add_slide(self.pptx.slide_layouts[6])
        pptx_text_box = pptx_slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(9),
            height=Inches(2.5)
        )
        pptx_text_frame = pptx_text_box.text_frame
        pptx_text_frame.word_wrap = True
        self.current_paragraph = pptx_text_frame.add_paragraph()

    def copy_paragraph_from_docx(self, docx_paragraph):
        for docx_run in docx_paragraph.runs:
            pptx_run = self.current_paragraph.add_run()
            pptx_run.text = docx_run.text
            pptx_run.font.bold = docx_run.bold
            pptx_run.font.italic = docx_run.italic
            pptx_run.font.underline = docx_run.underline
            pptx_run.font.name = docx_run.font.name
            font_color = str(docx_run.font.color.rgb)
            pptx_run.font.size = docx_run.font.size
            pptx_run.style = docx_run.style
            if font_color not in 'None':
                pptx_run.font.color.rgb = RGBColor(
                    r=int(font_color[0:2], 16),
                    g=int(font_color[2:4], 16),
                    b=int(font_color[4:6], 16)
                )
            if docx_run.font.highlight_color:
                hex_val = self.WD_COLOR_TO_RGB_STR[docx_run.font.highlight_color]
                self._set_highlight(pptx_run, hex_val)

    def save(self, path):
        self.pptx.save(path)

    # https://stackoverflow.com/a/62841826
    @staticmethod
    def _set_highlight(run, color):
        # get run properties
        rPr = run._r.get_or_add_rPr()
        # Create highlight element
        hl = OxmlElement("a:highlight")
        # Create specify RGB Colour element with color specified
        srgbClr = OxmlElement("a:srgbClr")
        setattr(srgbClr, "val", color)
        # Add colour specification to highlight element
        hl.append(srgbClr)
        # Add highlight element to run properties
        rPr.append(hl)
        return run


def main():
    # Ged rid off blank Tk window
    root = Tk()
    root.withdraw()

    # Load input docx file selected by user
    input_file = filedialog.askopenfilename(
        title='Select input DOCX file',
        filetypes=[('Word document', '.docx')],
    )
    input_file_path = Path(input_file)
    input_docx = Document(input_file_path)
    output_pptx = Pptx()

    output_pptx.new_blank_slide()
    for docx_paragraph in input_docx.paragraphs:
        if len(output_pptx.current_paragraph.text) != 0:
            # Don't create new blank slide for empty paragraph
            output_pptx.new_blank_slide()
        output_pptx.copy_paragraph_from_docx(docx_paragraph)

    # Save to user selected location with tkinter
    selected_output = filedialog.askdirectory(title='Select output folder')
    output_path = Path(selected_output).joinpath(input_file_path.name)
    output_path = output_path.with_suffix('.pptx')
    print(output_path)
    output_pptx.save(output_path)


if __name__ == '__main__':
    main()
